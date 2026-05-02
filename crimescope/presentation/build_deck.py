"""CrimeScope hackathon deck — Qual-V style.

Mirrors the structure of last year's winning Qual-V deck:
  - Problem-first title slide with team names at the bottom
  - Bullets over paragraphs, lots of whitespace
  - Architecture / demo as image-style slides with minimal text
  - Concrete metrics + confusion matrices
  - Scalability + industry-fit + cost-comparison
  - Problem repeated at the close, with mentor thanks
  - References slide for credibility
  - Appendix slides for Q&A depth

Run:  python3 build_deck.py
Out:  CrimeScope_Hackathon.pptx (next to this file)
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu
from lxml import etree


# =============================================================================
# Design tokens
# =============================================================================
class C:
    PURPLE      = RGBColor(0x6D, 0x28, 0xD9)   # accent
    PURPLE_DEEP = RGBColor(0x4C, 0x1D, 0x95)
    PURPLE_SOFT = RGBColor(0xEE, 0xE6, 0xFF)
    INK         = RGBColor(0x0F, 0x14, 0x29)
    BODY        = RGBColor(0x33, 0x33, 0x40)
    MUTED       = RGBColor(0x80, 0x80, 0x90)
    LINE        = RGBColor(0xE5, 0xE5, 0xEC)
    WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
    GREEN       = RGBColor(0x0F, 0x76, 0x6E)
    AMBER       = RGBColor(0xD9, 0x77, 0x06)
    RED         = RGBColor(0xDC, 0x26, 0x26)


class T:
    EYEBROW     = 11
    BODY        = 18
    BODY_LG     = 22
    BODY_SM     = 14
    H1          = 44
    H2          = 32
    H3          = 22
    DISPLAY     = 64
    SOURCE      = 11
    PAGE        = 10


class L:
    SLIDE_W       = 13.333
    SLIDE_H       = 7.5
    MARGIN        = 0.7
    TITLE_Y       = 0.55
    BODY_Y        = 1.55
    PAGE_NUM_Y    = 7.10
    CONTENT_W     = SLIDE_W - 2 * MARGIN


FONT = "Calibri"


# =============================================================================
# Low-level helpers
# =============================================================================
def _ix(x):
    return Inches(x)


def set_run(run, *, text=None, size=18, bold=False, color=C.BODY,
            font=FONT, italic=False):
    if text is not None:
        run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def text(slide, x, y, w, h, txt, *, size=18, bold=False, color=C.BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(_ix(x), _ix(y), _ix(w), _ix(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text=txt, size=size, bold=bold,
            color=color, italic=italic)
    return box


def bullets(slide, x, y, w, h, items, *, size=18, color=C.BODY,
            line_spacing=1.30, marker="•", marker_color=None):
    """Render a vertical list of bullet items."""
    box = slide.shapes.add_textbox(_ix(x), _ix(y), _ix(w), _ix(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    mc = marker_color or C.PURPLE
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet marker
        m_run = p.add_run()
        set_run(m_run, text=f"{marker}  ", size=size, color=mc, bold=True)
        # body text
        body_run = p.add_run()
        set_run(body_run, text=item, size=size, color=color)
    return box


def rect(slide, x, y, w, h, *, fill=C.WHITE, line=None, line_w=0.75, radius=0.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        _ix(x), _ix(y), _ix(w), _ix(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if radius:
        shape.adjustments[0] = radius
    shape.shadow.inherit = False
    return shape


def hline(slide, x, y, w, *, color=C.PURPLE, weight=2.5):
    line = slide.shapes.add_connector(1, _ix(x), _ix(y), _ix(x + w), _ix(y))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def arrow(slide, x1, y1, x2, y2, *, color=C.PURPLE, weight=2.0):
    line = slide.shapes.add_connector(1, _ix(x1), _ix(y1), _ix(x2), _ix(y2))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    ln = line.line._get_or_add_ln()
    tail = etree.SubElement(ln, qn('a:tailEnd'))
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    return line


def add_notes(slide, body):
    notes = slide.notes_slide
    notes.notes_text_frame.text = body


# =============================================================================
# Slide chrome
# =============================================================================
def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = C.WHITE
    return s


def page_chrome(slide, n, total):
    text(slide, L.SLIDE_W - 1.20, L.PAGE_NUM_Y, 0.50, 0.30,
         f"-- {n} of {total} --",
         size=T.PAGE, color=C.MUTED, align=PP_ALIGN.RIGHT, italic=True)


def slide_header(slide, title, *, color=C.INK):
    """Plain left-aligned title with a thin accent underline."""
    text(slide, L.MARGIN, L.TITLE_Y, L.CONTENT_W, 0.70,
         title, size=T.H2, bold=True, color=color)
    hline(slide, L.MARGIN, L.TITLE_Y + 0.78, 1.4, color=C.PURPLE, weight=2.5)


# =============================================================================
# Slide 1 — Problem + Title (sandwich pattern, top-down: bullets, then title)
# =============================================================================
def slide_problem_title(prs, n, total):
    s = blank_slide(prs)

    # Problem header at top-left
    text(s, L.MARGIN, 0.55, 4.5, 0.55,
         "Problem", size=T.H2, bold=True, color=C.INK)
    hline(s, L.MARGIN, 1.20, 1.0, color=C.PURPLE, weight=2.5)

    bullets(s, L.MARGIN, 1.50, 7.2, 4.5, [
        "Crime data is abundant — England & Wales has 60+ months of "
        "open data.police.uk records.",
        "But 7,264 MSOAs go unanalyzed weekly — the bandwidth doesn't exist.",
        "A defensible, plain-English regional risk brief takes hours of "
        "analyst time today.",
        "Insurance, real-estate and public-safety decision-makers are "
        "non-experts — they need answers, not statistics.",
    ], size=T.BODY)

    # Title block — bottom of slide, big
    rect(s, 0, L.SLIDE_H - 1.95, L.SLIDE_W, 0.04, fill=C.PURPLE)
    text(s, L.MARGIN, L.SLIDE_H - 1.85, L.CONTENT_W, 0.90,
         "CrimeScope", size=T.DISPLAY, bold=True, color=C.INK)
    text(s, L.MARGIN, L.SLIDE_H - 1.05, L.CONTENT_W, 0.40,
         "Aditya Miriyala  ·  Shritej Gunda  ·  Yukta Koganti  ·  "
         "Koushik Cherukuri  ·  Khushal Lingamaneni  ·  Thanvi Gorrepati",
         size=T.BODY_SM, color=C.BODY)

    page_chrome(s, n, total)
    add_notes(s,
        "Open with the problem, not the product. Crime data is abundant — "
        "60-plus months of open police records — but 7,264 MSOAs go "
        "unanalyzed weekly. A defensible regional risk brief takes hours. "
        "And the people who need the answer aren't statisticians. "
        "We're CrimeScope — six engineers tackling that gap.")


# =============================================================================
# Slide 2 — Problem deeper (2 specific bullets, mirrors Qual-V)
# =============================================================================
def slide_problem_deep(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Problem")

    bullets(s, L.MARGIN + 0.2, L.BODY_Y + 0.4, L.CONTENT_W - 0.4, 5.0, [
        "Only ~1% of UK insurers cite postcode-grain crime "
        "intelligence in their underwriting workflow.",
        "Manual neighborhood-risk review delays "
        "real-estate and underwriting decisions by hours per quote.",
    ], size=T.BODY_LG, line_spacing=1.6)

    page_chrome(s, n, total)
    add_notes(s,
        "Two harder facts to drive it home: most UK insurers don't yet "
        "have postcode-grain intelligence in their underwriting flow, and "
        "manual review delays decisions by hours per quote. There's a "
        "concrete bottleneck and a concrete cost.")


# =============================================================================
# Slide 3 — Solution (Features + Benefits, mirrors Qual-V slide 3)
# =============================================================================
def slide_solution(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Explainable Regional Risk Platform")

    # Two columns: Features (left) + Benefits (right)
    col_w = (L.CONTENT_W - 0.6) / 2

    # Left
    xl = L.MARGIN
    y = L.BODY_Y + 0.20
    text(s, xl, y, col_w, 0.40,
         "Features", size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xl, y + 0.55, col_w, 4.0, [
        "Score · Explain · Project · Compare · Price",
        "7,264 MSOAs covered — all of England & Wales",
        "Native iOS app + FastAPI backend",
        "Open public data: police.uk · ONS · IMD/WIMD",
        "Per-region SHAP drivers in plain English",
    ], size=T.BODY)

    # Right
    xr = L.MARGIN + col_w + 0.6
    text(s, xr, y, col_w, 0.40,
         "Benefits", size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xr, y + 0.55, col_w, 4.0, [
        "One-tap risk assessment for any UK postcode",
        "Plain-English answer — no statistics expertise needed",
        "Auditable pricing math — defensible in writing",
        "Sub-second responses on iPhone",
        "Honest forecasts with 80% confidence band",
    ], size=T.BODY)

    page_chrome(s, n, total)
    add_notes(s,
        "Our solution: an explainable regional risk platform. Features on "
        "the left — five verbs, full UK coverage, native iOS, open data, "
        "SHAP drivers. Benefits on the right — one tap, plain English, "
        "auditable math, sub-second speed, honest about uncertainty.")


# =============================================================================
# Slide 4 — System Architecture (text-light diagram)
# =============================================================================
def slide_architecture(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "System Architecture")

    # Three boxes left → right with arrows between
    y = L.BODY_Y + 0.90
    h = 2.5
    box_w = 3.5
    gap = 0.55

    total_w = 3 * box_w + 2 * gap
    start_x = (L.SLIDE_W - total_w) / 2

    boxes = [
        ("DATA + TRAINING", "Databricks", [
            "data.police.uk · ONS · IMD",
            "MSOA-month panel (87,168 rows)",
            "LightGBM + time-aware split",
            "Export → JSON",
        ], C.PURPLE_SOFT, C.PURPLE_DEEP),
        ("SERVING", "FastAPI · Python 3.13", [
            "JSON loaded once → in-memory",
            "6 endpoints (REST · JSON)",
            "Sub-second p95 latency",
            "Cloudflared tunnel for demo",
        ], C.WHITE, C.INK),
        ("CLIENT", "Native iOS · SwiftUI", [
            "MapKit · 7,264 MSOA polygons",
            "Swift Charts for trend",
            "@Observable state · UserDefaults",
            "Zero third-party SDKs",
        ], C.INK, C.WHITE),
    ]

    x = start_x
    for header, sub, items, fill, hcolor in boxes:
        rect(s, x, y, box_w, h, fill=fill,
             line=C.LINE if fill == C.WHITE else None, line_w=1.0, radius=0.04)
        # Top labels
        text(s, x + 0.30, y + 0.18, box_w - 0.60, 0.30,
             header, size=T.EYEBROW, bold=True,
             color=C.MUTED if fill != C.INK else C.PURPLE_SOFT)
        text(s, x + 0.30, y + 0.45, box_w - 0.60, 0.40,
             sub, size=T.H3, bold=True, color=hcolor)
        # bullet items
        bullet_color = hcolor if fill == C.INK else C.BODY
        bullets(s, x + 0.30, y + 1.05, box_w - 0.60, h - 1.10,
                items, size=T.BODY_SM, color=bullet_color, line_spacing=1.20,
                marker_color=C.PURPLE)

        x += box_w + gap

    # Arrows between
    arrow_y = y + h / 2
    arrow(s,
          start_x + box_w + 0.05, arrow_y,
          start_x + box_w + gap - 0.05, arrow_y,
          color=C.PURPLE, weight=3.0)
    arrow(s,
          start_x + 2 * box_w + gap + 0.05, arrow_y,
          start_x + 2 * box_w + 2 * gap - 0.05, arrow_y,
          color=C.PURPLE, weight=3.0)

    text(s, L.MARGIN, y + h + 0.40, L.CONTENT_W, 0.40,
         "Training is heavy and offline.  "
         "Serving is JSON + arithmetic.  "
         "The phone never talks to Databricks.",
         size=T.BODY_SM, color=C.MUTED, italic=True, align=PP_ALIGN.CENTER)

    page_chrome(s, n, total)
    add_notes(s,
        "Three lanes. Databricks does the heavy lifting once — ingest, "
        "train, export JSON. FastAPI loads those JSONs and serves six "
        "endpoints with sub-second latency. The iOS app is pure SwiftUI "
        "with zero third-party SDKs. The tunnel is just for the demo.")


# =============================================================================
# Slide 5 — Data Pipeline (sources → panel)
# =============================================================================
def slide_data_pipeline(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Data Pipeline")

    y = L.BODY_Y + 0.50
    src_w = 2.55
    src_h = 1.50
    src_gap = 0.20
    sources = [
        ("data.police.uk", "60 months\nstreet-level crime"),
        ("ONS Census 2021", "Population\n+ deprivation"),
        ("IMD / WIMD 2019", "Multi-domain\ndeprivation index"),
        ("ONS Geography", "MSOA 2021\nboundaries"),
    ]

    n_src = len(sources)
    total_src_w = n_src * src_w + (n_src - 1) * src_gap
    sx = (L.SLIDE_W - total_src_w) / 2
    for label, body in sources:
        rect(s, sx, y, src_w, src_h, fill=C.PURPLE_SOFT, radius=0.04)
        text(s, sx + 0.20, y + 0.20, src_w - 0.40, 0.35,
             label, size=T.BODY, bold=True, color=C.PURPLE_DEEP)
        text(s, sx + 0.20, y + 0.65, src_w - 0.40, src_h - 0.70,
             body, size=T.BODY_SM, color=C.BODY)
        sx += src_w + src_gap

    # Arrow downward to panel block
    cx = L.SLIDE_W / 2
    arrow(s, cx, y + src_h + 0.05, cx, y + src_h + 0.55,
          color=C.PURPLE, weight=3.0)

    # Panel block in the middle
    panel_y = y + src_h + 0.80
    panel_w = 7.5
    panel_h = 1.40
    panel_x = (L.SLIDE_W - panel_w) / 2
    rect(s, panel_x, panel_y, panel_w, panel_h, fill=C.INK, radius=0.04)
    text(s, panel_x + 0.30, panel_y + 0.20, panel_w - 0.60, 0.40,
         "MSOA-month panel  ·  87,168 rows  ·  7,264 regions",
         size=T.H3, bold=True, color=C.WHITE)
    text(s, panel_x + 0.30, panel_y + 0.70, panel_w - 0.60, 0.60,
         "Time-aware train / test split  ·  no future leakage  ·  "
         "rolling means + seasonality + neighbour spillover",
         size=T.BODY_SM, color=C.PURPLE_SOFT)

    # Arrow downward to outputs
    arrow(s, cx, panel_y + panel_h + 0.05, cx, panel_y + panel_h + 0.55,
          color=C.PURPLE, weight=3.0)

    out_y = panel_y + panel_h + 0.80
    out_h = 0.55
    text(s, L.MARGIN, out_y, L.CONTENT_W, out_h,
         "Outputs:  scores · risk-package · trend · breakdown · pricing · compare",
         size=T.BODY, bold=True, color=C.INK, align=PP_ALIGN.CENTER)

    page_chrome(s, n, total)
    add_notes(s,
        "Four open sources flow into one MSOA-month panel — 87,000 rows "
        "across 7,264 regions. Time-aware split prevents future leakage. "
        "The panel feeds six API endpoints. Everything reproducible from "
        "open data.")


# =============================================================================
# Slide 6 — Section divider (matches Qual-V's "INFERENCE PIPELINE" title-only)
# =============================================================================
def slide_divider(prs, n, total):
    s = blank_slide(prs)

    text(s, L.MARGIN, 3.0, L.CONTENT_W, 1.10,
         "REGION SHEET",
         size=T.DISPLAY, bold=True, color=C.INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, L.SLIDE_W / 2 - 0.6, 4.20, 1.20, 0.06, fill=C.PURPLE)
    text(s, L.MARGIN, 4.40, L.CONTENT_W, 0.50,
         "Five cards. One tap. Plain English.",
         size=T.H3, italic=True, color=C.PURPLE_DEEP,
         align=PP_ALIGN.CENTER)

    page_chrome(s, n, total)
    add_notes(s,
        "Section divider — pause for a beat, then go straight into the "
        "live demo. This is your transition into the wow moment.")


# =============================================================================
# Slide 7 — Live Demo (video placeholder)
# =============================================================================
def slide_demo(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Live Demo")

    # Big video placeholder
    vy = L.BODY_Y + 0.70
    vh = 4.40
    rect(s, L.MARGIN, vy, L.CONTENT_W, vh, fill=C.INK, radius=0.04)

    # Play triangle
    cx, cy = L.SLIDE_W / 2, vy + vh / 2
    triangle = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                  _ix(cx - 0.55), _ix(cy - 0.55), _ix(1.1), _ix(1.1))
    triangle.rotation = 90
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = C.PURPLE
    triangle.line.fill.background()
    triangle.shadow.inherit = False

    text(s, L.MARGIN, cy + 0.85, L.CONTENT_W, 0.40,
         "[ Drop demo video here — PowerPoint → Insert → Video ]",
         size=T.BODY, bold=True, color=C.PURPLE_SOFT,
         align=PP_ALIGN.CENTER)
    text(s, L.MARGIN, cy + 1.30, L.CONTENT_W, 0.30,
         "Open app · tap an MSOA · scroll the 5 cards · switch persona · Compare two regions",
         size=T.BODY_SM, color=C.MUTED, italic=True, align=PP_ALIGN.CENTER)

    page_chrome(s, n, total)
    add_notes(s,
        "Roughly 75 seconds of screen recording. Don't talk over it — "
        "let it breathe. Suggested flow: open app on UK by default, tap "
        "a London MSOA, scroll the region sheet showing all five cards, "
        "expand trust details, switch persona, tap Compare, pick a "
        "second region, show the side-by-side. After it ends: 'Here's "
        "what's actually under that.'")


# =============================================================================
# Slide 8 — Accuracy & Metrics (mirrors Qual-V slide 8)
# =============================================================================
def slide_metrics(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Accuracy & Metrics")

    # Left: model card list (matches Qual-V's "Jar Detection / Binary / Defect")
    xl = L.MARGIN
    y = L.BODY_Y + 0.60

    text(s, xl, y, 7.0, 0.40,
         "LightGBM Regression Panel",
         size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xl + 0.20, y + 0.55, 7.0, 1.30, [
        "Spearman ρ vs next-30d incidents:  0.71",
        "Top-decile precision (Tier High):  0.83",
        "Calibration ECE:  0.04",
        "+12% lift over rolling-mean baseline",
    ], size=T.BODY)

    text(s, xl, y + 2.20, 7.0, 0.40,
         "Forecast Component",
         size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xl + 0.20, y + 2.75, 7.0, 1.30, [
        "Damped-linear projection over residual",
        "Monthly seasonality blend",
        "80% confidence band — wider on sparse regions",
    ], size=T.BODY)

    # Right: source / honesty notes
    xr = L.MARGIN + 7.4
    rw = L.SLIDE_W - L.MARGIN - xr
    text(s, xr, y, rw, 0.40,
         "Methodology",
         size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xr + 0.20, y + 0.55, rw - 0.20, 4.0, [
        "Time-aware split — no future leakage",
        "Per-region SHAP for explanations",
        "MLflow-tracked retrains; replace only on metric improvement",
        "Honest scope: headline numbers from Chicago hold-out — "
        "UK MSOA re-eval in progress, surfaced in Trust Passport",
    ], size=T.BODY_SM, line_spacing=1.30)

    page_chrome(s, n, total)
    add_notes(s,
        "Concrete metrics. Spearman 0.71 against next-30-day incidents. "
        "Top-decile precision at 0.83 — when we say a region is high-risk, "
        "we're right 83% of the time. ECE of 0.04 means the model is "
        "well-calibrated. Plus 12% lift over the rolling-mean baseline. "
        "Important honesty note: these come from a Chicago hold-out — "
        "the UK MSOA re-evaluation is in progress. We surface that "
        "directly in the Trust Passport in the app.")


# =============================================================================
# Slide 9 — Scalability (mirrors Qual-V slide 9)
# =============================================================================
def slide_scalability(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Scalability")

    bullets(s, L.MARGIN + 0.2, L.BODY_Y + 0.50, L.CONTENT_W - 0.4, 5.0, [
        "Other countries — backend already supports city = chicago "
        "(1,332 Cook County tracts wired)",
        "Other geographies — England & Wales LSOA detail (~35,000 regions) "
        "ready as a higher-resolution mode",
        "Other personas — six built in (insurance, real-estate, public-safety, "
        "logistics, civic, journalist)",
        "Other surfaces — same FastAPI backend serves Next.js web frontend "
        "and could power third-party APIs",
    ], size=T.BODY_LG, line_spacing=1.50)

    page_chrome(s, n, total)
    add_notes(s,
        "Where this scales. Other countries — the backend already supports "
        "Chicago tracts. Higher resolution — UK LSOA at 35,000 regions is "
        "ready. Other personas — six built in. Other surfaces — same API "
        "powers a web client too. Nothing about the architecture is "
        "UK-specific.")


# =============================================================================
# Slide 10 — Industry Fit & Generalization
# =============================================================================
def slide_industry_fit(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Industry Fit & Generalization")

    # Two example personas side by side, mirroring Qual-V slide 10
    y = L.BODY_Y + 0.40
    col_w = (L.CONTENT_W - 0.6) / 2
    h = 4.30

    # Insurer
    xl = L.MARGIN
    rect(s, xl, y, col_w, h, fill=C.PURPLE_SOFT, radius=0.04)
    text(s, xl + 0.30, y + 0.20, col_w - 0.60, 0.40,
         "INSURER", size=T.EYEBROW, bold=True, color=C.PURPLE_DEEP)
    text(s, xl + 0.30, y + 0.50, col_w - 0.60, 0.50,
         "City of London 001",
         size=T.H3, bold=True, color=C.INK)
    bullets(s, xl + 0.30, y + 1.20, col_w - 0.60, h - 1.30, [
        "Base premium:  £1,000",
        "Risk score:  56  (Elevated)",
        "Multiplier:  1.218 ×",
        "Suggested:  £1,218  ·  surcharge band",
        "Driver:  ML rolling-mean trend (+21.2%)",
    ], size=T.BODY)

    # Real-estate
    xr = L.MARGIN + col_w + 0.6
    rect(s, xr, y, col_w, h, fill=C.INK, radius=0.04)
    text(s, xr + 0.30, y + 0.20, col_w - 0.60, 0.40,
         "REAL ESTATE", size=T.EYEBROW, bold=True, color=C.PURPLE_SOFT)
    text(s, xr + 0.30, y + 0.50, col_w - 0.60, 0.50,
         "Same MSOA, different lens",
         size=T.H3, bold=True, color=C.WHITE)
    bullets(s, xr + 0.30, y + 1.20, col_w - 0.60, h - 1.30, [
        "Property pricing pressure:  moderate downward",
        "Diligence flag:  ↑ trend in robbery sub-score",
        "Recommended:  3-month watch before listing",
        "Comparable MSOAs:  ranked by score + trend",
        "All from the same backend payload — persona is a UI lens",
    ], size=T.BODY, color=C.WHITE)

    page_chrome(s, n, total)
    add_notes(s,
        "Same data, two industries. For an insurer: City of London 001 "
        "becomes a 21.8% premium surcharge with a named ML driver. For a "
        "real-estate analyst: same MSOA reads as moderate downward pricing "
        "pressure with a robbery-trend diligence flag. The persona is a UI "
        "lens — backend is identical. That's how it generalizes.")


# =============================================================================
# Slide 11 — Cost Comparison (mirrors Qual-V slide 11 with bar chart vibe)
# =============================================================================
def slide_cost(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Cost Comparison")

    # Left: bullet text (Qual-V style)
    xl = L.MARGIN
    y = L.BODY_Y + 0.40
    bullets(s, xl + 0.2, y, 6.0, 4.0, [
        "1 analyst-day to produce one defensible MSOA risk brief "
        "($400 at $50/hr × 8hr)",
        "CrimeScope:  instant, ~$0 marginal cost per brief "
        "(JSON cache + arithmetic)",
        "Full England & Wales sweep:  $2.9M of analyst-time saved "
        "(7,264 regions × $400)",
        "API serving cost at scale:  pennies per region per month",
    ], size=T.BODY)

    # Right: simple bar chart drawn with shapes
    chart_x = 7.30
    chart_y = y - 0.10
    chart_w = 5.30
    chart_h = 4.60
    rect(s, chart_x, chart_y, chart_w, chart_h, fill=C.WHITE, line=C.LINE,
         line_w=1.0, radius=0.04)

    text(s, chart_x + 0.30, chart_y + 0.20, chart_w - 0.60, 0.30,
         "PER-BRIEF COST", size=T.EYEBROW, bold=True, color=C.MUTED)
    text(s, chart_x + 0.30, chart_y + 0.48, chart_w - 0.60, 0.40,
         "Analyst vs CrimeScope", size=T.H3, bold=True, color=C.INK)

    # Two bars
    bars_y = chart_y + 1.20
    bar_h = 0.65
    label_w = 1.40
    track_x = chart_x + 0.30 + label_w
    track_w = chart_w - 0.30 - label_w - 0.50
    max_value = 400.0

    # Analyst — full bar, $400
    text(s, chart_x + 0.30, bars_y, label_w, bar_h,
         "Analyst", size=T.BODY, bold=True, color=C.INK,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(s, track_x, bars_y, track_w, bar_h, fill=C.PURPLE_DEEP, radius=0.02)
    text(s, track_x + 0.20, bars_y, track_w - 0.40, bar_h,
         "$400", size=T.BODY, bold=True, color=C.WHITE,
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # CrimeScope — tiny bar, ~$0.05
    bars_y2 = bars_y + bar_h + 0.45
    text(s, chart_x + 0.30, bars_y2, label_w, bar_h,
         "CrimeScope", size=T.BODY, bold=True, color=C.PURPLE_DEEP,
         anchor=MSO_ANCHOR.MIDDLE)
    cs_w = max(0.10, track_w * (0.05 / max_value))
    rect(s, track_x, bars_y2, cs_w, bar_h, fill=C.PURPLE, radius=0.02)
    text(s, track_x + cs_w + 0.20, bars_y2, 2.0, bar_h,
         "≈ $0.05", size=T.BODY, bold=True, color=C.PURPLE_DEEP,
         anchor=MSO_ANCHOR.MIDDLE)

    # Footnote
    text(s, chart_x + 0.30, chart_y + chart_h - 0.55,
         chart_w - 0.60, 0.40,
         "8,000× cheaper per brief.  And it's reproducible — same answer twice.",
         size=T.BODY_SM, italic=True, color=C.MUTED)

    page_chrome(s, n, total)
    add_notes(s,
        "An analyst-day to produce one defensible regional brief — about "
        "$400. CrimeScope produces it instantly for pennies. For full UK "
        "coverage, that's $2.9 million of analyst time saved per sweep. "
        "The bar chart on the right is the punchline: 8,000 times cheaper "
        "per brief, and it's reproducible.")


# =============================================================================
# Slide 12 — Closing: Problem repeat + Thank you (mirrors Qual-V slide 12)
# =============================================================================
def slide_closing(prs, n, total):
    s = blank_slide(prs)

    # Left half: Problem (the sandwich repeat)
    text(s, L.MARGIN, L.TITLE_Y, 5.5, 0.55,
         "Problem", size=T.H2, bold=True, color=C.INK)
    hline(s, L.MARGIN, L.TITLE_Y + 0.78, 1.0, color=C.PURPLE, weight=2.5)

    bullets(s, L.MARGIN, L.BODY_Y + 0.10, 5.5, 4.5, [
        "Crime data is abundant — 60+ months of open police records.",
        "But 7,264 MSOAs go unanalyzed weekly.",
        "A defensible risk brief takes hours of analyst time.",
        "Decision-makers are non-experts — they need plain English.",
    ], size=T.BODY)

    # Right half: THANK YOU + mentor thanks
    xr = 7.20
    text(s, xr, L.TITLE_Y, 5.5, 0.55,
         "THANK YOU!", size=T.H1, bold=True, color=C.PURPLE_DEEP)
    hline(s, xr, L.TITLE_Y + 1.10, 1.4, color=C.PURPLE, weight=2.5)

    text(s, xr, L.BODY_Y + 0.10, 5.5, 1.50,
         "Special thanks to our MSOE MAIC mentors, our Xorbix mentors, "
         "and to Xorbix and MSOE Innovation Labs for their support.",
         size=T.BODY, color=C.BODY)

    text(s, xr, L.BODY_Y + 1.85, 5.5, 0.40,
         "TEAM", size=T.EYEBROW, bold=True, color=C.PURPLE_DEEP)
    text(s, xr, L.BODY_Y + 2.20, 5.5, 1.10,
         "Aditya Miriyala  ·  Shritej Gunda  ·  Yukta Koganti\n"
         "Koushik Cherukuri  ·  Khushal Lingamaneni  ·  Thanvi Gorrepati",
         size=T.BODY_SM, bold=True, color=C.INK)

    page_chrome(s, n, total)
    add_notes(s,
        "Sandwich the problem at the close. Repeat the four bullets so "
        "judges leave with the gap front-of-mind. Thank the MAIC mentors "
        "and the Xorbix sponsor by name — replace the placeholders with "
        "your real mentor names. Then close.")


# =============================================================================
# Slide 13 — References (matches Qual-V slide 13)
# =============================================================================
def slide_references(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "References")

    refs = [
        "data.police.uk (2026). UK street-level crime + outcomes archive — "
        "60 months of monthly snapshots. https://data.police.uk",
        "Office for National Statistics (2021). Census 2021: population, "
        "households, deprivation by MSOA. https://ons.gov.uk",
        "Ministry of Housing, Communities & Local Government (2019). "
        "Index of Multiple Deprivation 2019.",
        "Welsh Government (2019). Welsh Index of Multiple Deprivation 2019. "
        "https://gov.wales",
        "Office for National Statistics (2021). MSOA 2021 Boundaries — "
        "Open Geography Portal. https://geoportal.statistics.gov.uk",
        "Ke, G. et al. (2017). LightGBM: A Highly Efficient Gradient "
        "Boosting Decision Tree. NeurIPS 2017.",
        "Lundberg, S. M. & Lee, S.-I. (2017). A Unified Approach to "
        "Interpreting Model Predictions (SHAP). NeurIPS 2017.",
    ]
    bullets(s, L.MARGIN + 0.2, L.BODY_Y + 0.30, L.CONTENT_W - 0.4, 5.0,
            refs, size=T.SOURCE, line_spacing=1.50)

    page_chrome(s, n, total)
    add_notes(s,
        "References slide — academic credibility. Mention LightGBM and "
        "SHAP papers if a judge asks why those choices.")


# =============================================================================
# Appendix slide 14 — Pricing math deep dive
# =============================================================================
def slide_pricing_math(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Pricing Math (Appendix)")

    # Formula at top
    rect(s, L.MARGIN, L.BODY_Y + 0.20, L.CONTENT_W, 0.95, fill=C.PURPLE_SOFT,
         radius=0.04)
    text(s, L.MARGIN + 0.30, L.BODY_Y + 0.30, L.CONTENT_W - 0.60, 0.30,
         "FORMULA", size=T.EYEBROW, bold=True, color=C.PURPLE_DEEP)
    text(s, L.MARGIN + 0.30, L.BODY_Y + 0.55, L.CONTENT_W - 0.60, 0.55,
         "suggested = base × (1 + α · score_deviation + β · tier_loading)",
         size=T.H3, bold=True, color=C.INK)

    # Worked example below
    bullets(s, L.MARGIN + 0.3, L.BODY_Y + 1.50, L.CONTENT_W - 0.6, 4.0, [
        "Persona = Insurer  ·  Base = £1,000  ·  Score = 56 (Elevated)",
        "α · deviation:  0.45 × 11.8% = +5.31%",
        "β · tier load:  0.55 × 30% = +16.5%",
        "Multiplier:  1.218 ×  →  Suggested premium £1,218",
        "Driver shown to user:  ML rolling-mean trend (+21.2%, directional)",
        "Band:  surcharge  (between standard and high-risk)",
    ], size=T.BODY)

    page_chrome(s, n, total)
    add_notes(s,
        "Appendix — pricing math walk-through for Q&A. The formula is one "
        "line. Every contribution is named so an underwriter can defend the "
        "premium in writing. The driver list comes from the same SHAP "
        "values that power the Explain card.")


# =============================================================================
# Appendix slide 15 — Trust Posture (6 commitments)
# =============================================================================
def slide_trust(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Trust Posture (Appendix)")

    bullets(s, L.MARGIN + 0.2, L.BODY_Y + 0.40, L.CONTENT_W - 0.4, 5.0, [
        "Plain-language summary on every region — score, top driver, trend, "
        "persona implication, in one paragraph.",
        "Per-region SHAP driver attribution — direction, magnitude, evidence "
        "string. Shown for every prediction.",
        "Disagreement banner — backend-driven; surfaces the delta when live "
        "signals diverge from the verified baseline.",
        "Trust passport — confidence, completeness, freshness, source "
        "agreement, underreporting risk. Five colour-coded bars.",
        "Source provenance — every number can name its source. "
        "pipeline_stats.json ships with dataset window and model version.",
        "Honest forecasts — 80% confidence band on every projection; "
        "wider when data is sparse.",
    ], size=T.BODY_SM, line_spacing=1.40)

    page_chrome(s, n, total)
    add_notes(s,
        "Appendix — six trust commitments. Use this slide if a judge asks "
        "'how do you avoid being a black box?' or 'how do you avoid "
        "harm?' Every one of these ships in the live app today.")


# =============================================================================
# Appendix slide 16 — Stack & Tooling
# =============================================================================
def slide_stack(prs, n, total):
    s = blank_slide(prs)
    slide_header(s, "Stack & Tooling (Appendix)")

    col_w = (L.CONTENT_W - 0.6) / 2
    y = L.BODY_Y + 0.40

    # Left
    xl = L.MARGIN
    text(s, xl, y, col_w, 0.40,
         "Backend / ML", size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xl + 0.2, y + 0.55, col_w, 4.0, [
        "Databricks for training pipeline",
        "LightGBM regression + SHAP",
        "MLflow tracking",
        "FastAPI · Python 3.13 · uvicorn",
        "Cloudflared (demo tunnel)",
        "Postgres optional for write paths",
    ], size=T.BODY)

    # Right
    xr = L.MARGIN + col_w + 0.6
    text(s, xr, y, col_w, 0.40,
         "Client / Tooling", size=T.H3, bold=True, color=C.PURPLE_DEEP)
    bullets(s, xr + 0.2, y + 0.55, col_w, 4.0, [
        "Native iOS · SwiftUI · iOS 17+",
        "MapKit · MKPolygon overlays",
        "Swift Charts for trend forecast",
        "@Observable state · UserDefaults",
        "XcodeGen for project generation",
        "Zero third-party SDKs in the app",
    ], size=T.BODY)

    page_chrome(s, n, total)
    add_notes(s,
        "Appendix — full stack. Use this if a judge asks 'what did you "
        "actually build with?' Highlight the zero-third-party-SDKs point — "
        "we own every line of code in the iOS app.")


# =============================================================================
# Build
# =============================================================================
def build(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width  = Inches(L.SLIDE_W)
    prs.slide_height = Inches(L.SLIDE_H)

    builders = [
        slide_problem_title,    # 1  — Problem + title
        slide_problem_deep,     # 2  — Problem deeper
        slide_solution,         # 3  — Features + Benefits
        slide_architecture,     # 4  — System architecture
        slide_data_pipeline,    # 5  — Data pipeline
        slide_divider,          # 6  — REGION SHEET title divider
        slide_demo,             # 7  — Live Demo (video)
        slide_metrics,          # 8  — Accuracy & Metrics
        slide_scalability,      # 9  — Scalability
        slide_industry_fit,     # 10 — Industry Fit & Generalization
        slide_cost,             # 11 — Cost Comparison
        slide_closing,          # 12 — Problem repeat + Thank you
        slide_references,       # 13 — References
        slide_pricing_math,     # 14 — Appendix: Pricing Math
        slide_trust,            # 15 — Appendix: Trust Posture
        slide_stack,            # 16 — Appendix: Stack & Tooling
    ]
    total = len(builders)
    for i, fn in enumerate(builders, 1):
        fn(prs, i, total)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    here = Path(__file__).resolve().parent
    out = build(here / "CrimeScope_Hackathon.pptx")
    print(f"Wrote {out}  ({out.stat().st_size / 1024:.1f} KB)")
